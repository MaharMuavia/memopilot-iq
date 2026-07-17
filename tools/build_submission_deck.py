"""Build the editable MemoPilot IQ hackathon deck.

The deck marks cloud proof as pending. Rebuild it after the final Alibaba
deployment and replace placeholders with only verified URLs and evidence.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "assets" / "memopilot-iq-hackathon-deck.pptx"
NAVY, BLUE, SLATE, MUTED, WHITE = "0F172A", "2563EB", "475569", "64748B", "FFFFFF"
PALE_BLUE, PALE_INDIGO, PALE_AMBER = "EFF6FF", "EEF2FF", "FFFBEB"


def color(value: str):
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(value)


def text(slide, value: str, x: float, y: float, w: float, h: float, *, size: float = 18,
         fill: str = SLATE, bold: bool = False, align=PP_ALIGN.LEFT) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = value
    paragraph.alignment = align
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = color(fill)
    paragraph.font.bold = bold


def bullets(slide, items: list[str], x: float, y: float, w: float, h: float) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.bullet = True
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = color(SLATE)
        paragraph.space_after = Pt(10)


def card(slide, title: str, body: str, x: float, y: float, w: float, h: float, *, fill: str = PALE_BLUE) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.color.rgb = color("DBEAFE")
    text(slide, title, x + 0.2, y + 0.18, w - 0.4, 0.28, size=16, fill=NAVY, bold=True)
    text(slide, body, x + 0.2, y + 0.62, w - 0.4, h - 0.76, size=13)


def header(slide, title: str, number: int, subtitle: str = "") -> None:
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.78))
    band.fill.solid()
    band.fill.fore_color.rgb = color(NAVY)
    band.line.fill.background()
    text(slide, title, 0.55, 0.16, 9.5, 0.38, size=26, fill=WHITE, bold=True)
    text(slide, subtitle, 10.0, 0.23, 2.75, 0.2, size=10, fill="BFDBFE", align=PP_ALIGN.RIGHT)
    text(slide, "MemoPilot IQ · Qwen Cloud Global AI Hackathon · Track 1: MemoryAgent", 0.55, 7.1, 10.8, 0.18, size=9, fill=MUTED)
    text(slide, str(number), 12.2, 7.1, 0.5, 0.18, size=9, fill=MUTED, align=PP_ALIGN.RIGHT)


def slide(prs: Presentation, number: int, title: str, subtitle: str = ""):
    page = prs.slides.add_slide(prs.slide_layouts[6])
    page.background.fill.solid()
    page.background.fill.fore_color.rgb = color(WHITE)
    header(page, title, number, subtitle)
    return page


def build() -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    page = prs.slides.add_slide(prs.slide_layouts[6])
    page.background.fill.solid()
    page.background.fill.fore_color.rgb = color(NAVY)
    text(page, "MemoPilot IQ", 0.85, 1.45, 11.6, 0.7, size=42, fill=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(page, "Auditable memory governance for persistent AI agents", 1.1, 2.3, 11.1, 0.45, size=22, fill="BFDBFE", align=PP_ALIGN.CENTER)
    text(page, "Qwen Cloud Global AI Hackathon · Track 1: MemoryAgent", 1.1, 3.05, 11.1, 0.28, size=16, fill="C7D2FE", align=PP_ALIGN.CENTER)
    text(page, "Current build: React + Vite · FastAPI · Qwen via DashScope · SQLite locally", 1.1, 4.25, 11.1, 0.3, size=15, fill=WHITE, align=PP_ALIGN.CENTER)
    text(page, "Cloud deployment proof will be added only after verification.", 1.1, 4.75, 11.1, 0.24, size=12, fill="FDE68A", align=PP_ALIGN.CENTER)
    text(page, "1", 12.2, 7.1, 0.5, 0.18, size=9, fill="BFDBFE", align=PP_ALIGN.RIGHT)

    page = slide(prs, 2, "The problem", "Long conversations are not durable memory")
    card(page, "Context loss", "New sessions forget preferences, decisions, and constraints.", 0.7, 1.35, 3.8, 2.15)
    card(page, "Stale decisions", "Naive retrieval can replay choices the user later replaced.", 4.75, 1.35, 3.8, 2.15, fill=PALE_INDIGO)
    card(page, "Token pressure", "Full-history prompts waste context on irrelevant or expired material.", 8.8, 1.35, 3.8, 2.15)
    text(page, "The goal is not to store everything. It is to retain the right facts, retire stale ones, and explain every selection.", 1.25, 4.35, 10.8, 0.75, size=24, fill=NAVY, bold=True, align=PP_ALIGN.CENTER)

    page = slide(prs, 3, "The solution", "A lifecycle-aware memory layer")
    bullets(page, [
        "Extract typed records instead of keeping raw chat history.",
        "Retrieve with semantics, keywords, tags, project scope, and lifecycle state.",
        "Respect a hard 2,500-token memory budget and prioritise critical constraints.",
        "Supersede, expire, archive, and delete without losing the audit trail.",
        "Expose a Memory Trace for included and skipped memories.",
    ], 0.9, 1.35, 6.2, 4.9)
    card(page, "Judge-facing proof", "Cross-session recall\nSupersession without stale leakage\nTransparent trace accounting", 7.65, 1.65, 4.65, 2.55, fill=PALE_INDIGO)
    card(page, "Product truth", "Current implementation facts stay separate from user preferences and future plans.", 7.65, 4.55, 4.65, 1.35, fill=PALE_AMBER)

    page = slide(prs, 4, "Architecture", "Current implementation and cloud adapter boundary")
    nodes = ["React + Vite", "FastAPI", "MemoPilot\nMemory Layer", "Qwen Cloud"]
    for index, label in enumerate(nodes):
        x = 0.55 + index * 1.95
        card(page, label, "", x, 2.15, 1.55, 1.0, fill=PALE_INDIGO if index == 2 else PALE_BLUE)
        if index < len(nodes) - 1:
            text(page, "→", x + 1.56, 2.47, 0.34, 0.25, size=21, fill=BLUE, align=PP_ALIGN.CENTER)
    text(page, "SQLite + local vectors today · Alibaba Tablestore + OSS adapters after deployment", 0.75, 3.8, 7.2, 0.45, size=16, fill=SLATE, align=PP_ALIGN.CENTER)
    card(page, "Implemented now", "React 18 + Vite frontend\nFastAPI backend\nQwen chat, extraction, and embeddings through DashScope\nSQLite + local vector retrieval", 8.5, 1.25, 4.1, 2.35)
    card(page, "Deploy after the local gate", "Tablestore and OSS adapters exist. Do not label them live until health, persistence, and cloud evidence are verified.", 8.5, 4.0, 4.1, 1.75, fill=PALE_AMBER)

    page = slide(prs, 5, "Memory lifecycle", "Keep history without replaying stale decisions")
    stages = ["Create", "Retrieve", "Use", "Supersede", "Expire", "Audit"]
    for index, label in enumerate(stages):
        x = 0.65 + index * 2.08
        shape = page.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.2), Inches(1.55), Inches(0.86))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color(PALE_INDIGO if index in (3, 4) else PALE_BLUE)
        shape.line.color.rgb = color("C7D2FE")
        text(page, label, x, 2.47, 1.55, 0.22, size=16, fill=NAVY, bold=True, align=PP_ALIGN.CENTER)
        if index < len(stages) - 1:
            text(page, "→", x + 1.58, 2.42, 0.4, 0.3, size=22, fill=BLUE, align=PP_ALIGN.CENTER)
    text(page, "A planned migration can supersede an earlier preference while the assistant still states the verified stack in the submitted build today.", 1.25, 4.2, 10.8, 0.55, size=23, fill=NAVY, bold=True, align=PP_ALIGN.CENTER)
    text(page, "User intent is retained, but it does not rewrite repository facts.", 1.35, 5.05, 10.6, 0.35, size=15, fill=MUTED, align=PP_ALIGN.CENTER)

    page = slide(prs, 6, "Retrieval, safety, and control", "Transparent by design")
    card(page, "Retrieval", "Hybrid scoring combines relevance, importance, recency, usage, project match, criticality, and lifecycle penalties.", 0.7, 1.35, 3.85, 2.45)
    card(page, "Safety", "Secret patterns are redacted before persistence. Production supports server-side API keys and rate limiting.", 4.75, 1.35, 3.85, 2.45, fill=PALE_INDIGO)
    card(page, "User control", "Pin, archive, edit, delete, export, forget-all, timeline, graph, analytics, and history.", 8.8, 1.35, 3.85, 2.45)
    text(page, "The Memory Trace makes retrieval inspectable: included memory, skipped memory, scores, reasons, and token cost.", 1.05, 4.65, 11.2, 0.5, size=23, fill=NAVY, bold=True, align=PP_ALIGN.CENTER)

    page = slide(prs, 7, "Qwen Cloud integration", "Three live model roles")
    card(page, "Chat", "Qwen generates answers from the current message plus selected memory context.", 0.85, 1.45, 3.65, 2.2)
    card(page, "Extraction", "Qwen returns structured memory actions: create, update, supersede, or forget.", 4.85, 1.45, 3.65, 2.2, fill=PALE_INDIGO)
    card(page, "Embeddings", "Qwen embeddings support semantic retrieval alongside sparse and structured signals.", 8.85, 1.45, 3.65, 2.2)
    text(page, "Offline fallback keeps setup and tests reproducible. Final evidence must use the deployed Qwen-backed configuration.", 1.05, 4.75, 11.2, 0.48, size=20, align=PP_ALIGN.CENTER)

    page = slide(prs, 8, "Evaluation", "Evidence, not hard-coded claims")
    bullets(page, [
        "24-scenario diagnostic compares memory-augmented answers with a no-memory baseline.",
        "Reports strict answer checks, context recall, stale-memory avoidance, token reduction, and retrieval latency.",
        "The UI displays the report generated by the active configuration; it does not present preset winner metrics.",
        "Final evidence saves raw JSON with the deployed commit SHA, model, provider status, and UTC timestamp.",
    ], 0.95, 1.35, 7.0, 4.9)
    card(page, "Current status", "Reproducible local diagnostic: ready.\nFinal model-backed cloud benchmark: pending deployment.", 8.45, 2.0, 3.9, 1.85, fill=PALE_AMBER)

    page = slide(prs, 9, "Judge demo", "A clear three-minute story")
    bullets(page, [
        "Create project preferences and a critical no-secrets constraint.",
        "Open a new session and show cross-session recall in the Memory Trace.",
        "Record a planned future migration and show the supersession event.",
        "Ask for current versus planned stack: React + Vite today; Next.js only as the next-step plan.",
        "Show the evaluation report, then live cloud health and persistence proof after deployment.",
    ], 0.9, 1.25, 7.0, 5.2)
    card(page, "Recording rule", "Use the exact final deployed build. Keep the video public, English or subtitled, under three minutes, and free of credentials or copyrighted music.", 8.35, 2.0, 3.95, 2.25, fill=PALE_AMBER)

    page = slide(prs, 10, "Alibaba Cloud deployment proof", "Pending until the local release gate is green")
    card(page, "Cloud services", "DashScope/Qwen for model calls\nTablestore for persistent memory/events\nOSS for redacted artifacts\nECS, Function Compute, or ACK for hosting", 0.8, 1.3, 4.0, 3.2)
    card(page, "Required verification", "Public signed-out URL\n/health with Qwen + actual cloud mode\nMemory survives restart\nTablestore and OSS evidence\nRaw deployed benchmark report", 4.95, 1.3, 3.55, 3.2, fill=PALE_INDIGO)
    card(page, "No fabricated proof", "Replace this slide's pending status only after the exact public deployment, screenshots, and raw benchmark output are verified.", 8.7, 1.3, 3.8, 3.2, fill=PALE_AMBER)

    page = slide(prs, 11, "Submission checklist", "Finish only with verified evidence")
    bullets(page, [
        "Public repository: github.com/MaharMuavia/memopilot-iq",
        "Track 1: MemoryAgent, project description, architecture diagram, and MIT license.",
        "Public working URL and deployment proof from Alibaba Cloud.",
        "Public under-three-minute demo video and final deck with working links.",
        "Published build story and raw final evaluation report tied to the submitted commit.",
    ], 0.9, 1.25, 7.25, 5.2)
    card(page, "Before submitting", "Check every public link while signed out. Confirm the deployed app, README, video, deck, health response, benchmark, and Devpost description agree on the same final build.", 8.45, 2.0, 3.9, 2.3, fill=PALE_AMBER)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    build()
